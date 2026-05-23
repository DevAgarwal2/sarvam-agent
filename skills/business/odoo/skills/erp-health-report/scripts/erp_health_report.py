#!/usr/bin/env python3
"""
ERP Health Report — cross-skill orchestrator.

Queries all Odoo modules, generates matplotlib charts,
compiles a PPTX report, and optionally emails via Gmail.

Usage:
    python erp_health_report.py [--email to@example.com] [--output-dir ./report] [--json]
"""
import argparse
import json
import os
import subprocess
import sys
import time
from datetime import datetime
from pathlib import Path

HERMES_HOME = Path(os.environ.get("HERMES_HOME", Path.home() / ".hermes"))
ODOO_SCRIPT = HERMES_HOME / "skills" / "business" / "odoo" / "scripts" / "odoo_api.py"
GAPI_SCRIPT = HERMES_HOME / "skills" / "productivity" / "google-workspace" / "scripts" / "google_api.py"


def eprint(*args, **kwargs):
    print(*args, file=sys.stderr, **kwargs)


def run_odoo(*args):
    result = subprocess.run(
        [sys.executable, str(ODOO_SCRIPT)] + list(args),
        capture_output=True, text=True, timeout=60
    )
    if result.returncode != 0:
        return {"ok": False, "error": result.stderr.strip()}
    try:
        return json.loads(result.stdout)
    except json.JSONDecodeError:
        return {"ok": False, "error": f"Invalid JSON: {result.stdout[:200]}"}


def query_module_stats():
    modules = ["crm", "sales", "inventory", "manufacturing", "purchase", "accounting", "hr", "expenses", "contacts"]
    results = {}
    for mod in modules:
        r = run_odoo(mod, "statistics")
        if r.get("ok"):
            results[mod] = r.get("result", {})
        else:
            results[mod] = {"error": r.get("error", "unknown")}
    return results


def query_additional_data():
    data = {}
    r = run_odoo("check")
    if r.get("ok"):
        data["server"] = {
            "url": r.get("url"),
            "db": r.get("db"),
            "user": r.get("user"),
            "version": r.get("version", {}).get("server_version") if isinstance(r.get("version"), dict) else str(r.get("version", "")),
        }
    r = run_odoo("model", "product.product", "count", "--domain", "[]")
    if r.get("ok"):
        data["total_products"] = r["result"]
    r = run_odoo("model", "res.partner", "count", "--domain", "[]")
    if r.get("ok"):
        data["total_partners"] = r["result"]
    r = run_odoo("model", "product.product", "search", "--domain", '[["qty_available","<=",0]]', "--limit", "1", "--fields", "id")
    if r.get("ok"):
        data["out_of_stock_count"] = len(r.get("result", []))
    else:
        data["out_of_stock_count"] = 0
    return data


def generate_charts(stats, output_dir):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    charts_dir = output_dir / "charts"
    charts_dir.mkdir(parents=True, exist_ok=True)
    chart_files = {}

    color_map = {
        "crm": "#3498db",
        "sales": "#2ecc71",
        "inventory": "#e67e22",
        "manufacturing": "#9b59b6",
        "purchase": "#1abc9c",
        "accounting": "#e74c3c",
        "hr": "#f39c12",
        "expenses": "#95a5a6",
        "contacts": "#34495e",
    }

    for mod, data in stats.items():
        if not isinstance(data, dict) or "error" in data:
            continue
        fig, ax = plt.subplots(figsize=(8, 4.5))
        fig.patch.set_facecolor("#f8f9fa")
        ax.set_facecolor("#f8f9fa")

        items = list(data.items())
        if not items:
            continue

        labels = [str(k).replace("_", " ").title() for k, _ in items]
        values = [v if isinstance(v, (int, float)) else 0 for _, v in items]

        bars = ax.bar(range(len(items)), values, color=color_map.get(mod, "#3498db"), width=0.6, edgecolor="white", linewidth=1.2)
        for bar, v in zip(bars, values):
            ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + max(values) * 0.02,
                    str(v), ha="center", va="bottom", fontsize=10, fontweight="bold")

        ax.set_xticks(range(len(items)))
        ax.set_xticklabels(labels, fontsize=9, rotation=25, ha="right")
        ax.set_title(f"{mod.upper()} — Health Overview", fontsize=14, fontweight="bold", pad=12)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_color("#ddd")
        ax.spines["bottom"].set_color("#ddd")
        ax.tick_params(colors="#666")
        plt.tight_layout()

        path = charts_dir / f"{mod}.png"
        fig.savefig(path, dpi=150, bbox_inches="tight")
        plt.close(fig)
        chart_files[mod] = str(path)

    fig_combined, axes = plt.subplots(3, 3, figsize=(18, 12))
    fig_combined.patch.set_facecolor("#f8f9fa")
    fig_combined.suptitle("ERP Health Report — Module Overview", fontsize=20, fontweight="bold", y=0.98)

    for idx, (mod, data) in enumerate(stats.items()):
        if not isinstance(data, dict) or "error" in data:
            continue
        row, col = divmod(idx, 3)
        ax = axes[row, col]
        ax.set_facecolor("#f8f9fa")

        items = list(data.items())
        labels = [str(k).replace("_", " ").title() for k, _ in items][:6]
        values = [v if isinstance(v, (int, float)) else 0 for _, v in items][:6]

        ax.bar(range(len(labels)), values, color=color_map.get(mod, "#3498db"), width=0.6, edgecolor="white", linewidth=1.2)
        ax.set_xticks(range(len(labels)))
        ax.set_xticklabels(labels, fontsize=8, rotation=20, ha="right")
        ax.set_title(f"{mod.upper()}", fontsize=12, fontweight="bold")
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_color("#ddd")
        ax.spines["bottom"].set_color("#ddd")

    for i in range(len(stats), 9):
        row, col = divmod(i, 3)
        axes[row, col].axis("off")

    plt.tight_layout(rect=[0, 0, 1, 0.96])
    combined_path = charts_dir / "all-modules.png"
    fig_combined.savefig(combined_path, dpi=150, bbox_inches="tight")
    plt.close(fig_combined)
    chart_files["combined"] = str(combined_path)

    return chart_files


def generate_infographic_prompt(stats, server_info):
    now = datetime.now().strftime("%Y-%m-%d %H:%M")
    lines = []
    lines.append(f"# ERP Health Report — Infographic")
    lines.append(f"Generated: {now}")
    lines.append(f"Server: {server_info.get('server', {}).get('url', 'N/A')} / {server_info.get('server', {}).get('db', 'N/A')}")
    lines.append(f"Odoo Version: {server_info.get('server', {}).get('version', 'N/A')}")
    lines.append("")
    lines.append("## Summary Statistics")
    lines.append("")

    for mod, data in stats.items():
        if not isinstance(data, dict) or "error" in data:
            continue
        clean = {k: v for k, v in data.items() if isinstance(v, (int, float))}
        if clean:
            summary = " | ".join(f"{k}: {v}" for k, v in clean.items())
            lines.append(f"- **{mod.upper()}**: {summary}")

    lines.append("")
    lines.append("## Recommended Layout + Style")
    lines.append("- Layout: `dashboard` (best for metrics/KPIs)")
    lines.append("- Style: `corporate-memphis` (flat vector, vibrant — professional)")
    lines.append("- Aspect: landscape (16:9)")
    return "\n".join(lines)


def generate_pptx(stats, server_info, chart_files, output_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_slide_bg(slide, color=RGBColor(0xF8, 0xF9, 0xFA)):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=RGBColor(0x33, 0x33, 0x33), align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        return tf

    def add_table(slide, left, top, width, height, headers, rows):
        table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
        table = table_shape.table

        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)

        for r, row_data in enumerate(rows):
            for c, val in enumerate(row_data):
                cell = table.cell(r + 1, c)
                cell.text = str(val)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF3, 0xF4)

        return table

    add_slide_bg(prs.slides.add_slide(prs.slide_layouts[6]))
    slide = prs.slides[-1]

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(3.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    shape.line.fill.background()

    add_textbox(slide, 0.8, 0.6, 11, 1.2, "ERP Health Report", size=40, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_textbox(slide, 0.8, 1.8, 11, 0.8, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}", size=16, color=RGBColor(0xEC, 0xF0, 0xF1))

    si = server_info.get("server", {})
    details = [
        f"Odoo: {si.get('version', 'N/A')}",
        f"Database: {si.get('db', 'N/A')}",
        f"User: {si.get('user', 'N/A')}",
        f"Total Products: {server_info.get('total_products', 'N/A')}",
        f"Total Partners: {server_info.get('total_partners', 'N/A')}",
    ]
    for i, d in enumerate(details):
        add_textbox(slide, 0.8 + i * 2.5, 3.6, 2.3, 0.5, d, size=11, color=RGBColor(0x55, 0x55, 0x55))

    def _hex_to_rgb(h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    color_map = {
        "crm": ("CRM", "#3498db"), "sales": ("Sales", "#2ecc71"), "inventory": ("Inventory", "#e67e22"),
        "manufacturing": ("Manufacturing", "#9b59b6"), "purchase": ("Purchase", "#1abc9c"),
        "accounting": ("Accounting", "#e74c3c"), "hr": ("HR", "#f39c12"),
        "expenses": ("Expenses", "#95a5a6"), "contacts": ("Contacts", "#34495e"),
    }

    for idx, (mod, (label, color_hex)) in enumerate(color_map.items()):
        data = stats.get(mod, {})
        if not isinstance(data, dict) or "error" in data:
            continue

        x = 0.5 + (idx % 3) * 4.3
        y = 4.6 + (idx // 3) * 1.0

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(0.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(color_hex)
        shape.line.fill.background()

        vals = ", ".join(f"{k}={v}" for k, v in data.items() if isinstance(v, (int, float)))
        add_textbox(slide, x + 0.15, y + 0.05, 3.5, 0.6, f"{label}: {vals}", size=9, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    for mod, (label, _) in color_map.items():
        chart_path = chart_files.get(mod)
        if chart_path and Path(chart_path).exists():
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_slide_bg(slide)
            add_textbox(slide, 0.5, 0.2, 12, 0.6, f"{label} — Detailed Metrics", size=24, bold=True, color=RGBColor(0x2C, 0x3E, 0x50))
            slide.shapes.add_picture(chart_path, Inches(1.5), Inches(1.0), Inches(10), Inches(5.5))

            data = stats.get(mod, {})
            if isinstance(data, dict) and "error" not in data:
                headers = ["Metric", "Value"]
                rows = [[k.replace("_", " ").title(), str(v)] for k, v in data.items() if isinstance(v, (int, float))]
                if rows:
                    add_table(slide, 0.5, 6.4, 5, 0.8, headers, rows)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_textbox(slide, 0.5, 0.2, 12, 0.6, "Combined Overview", size=24, bold=True, color=RGBColor(0x2C, 0x3E, 0x50))
    combined = chart_files.get("combined")
    if combined and Path(combined).exists():
        slide.shapes.add_picture(combined, Inches(0.3), Inches(0.9), Inches(12.7), Inches(6.5))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_textbox(slide, 0.5, 0.2, 12, 0.6, "Key Recommendations", size=24, bold=True, color=RGBColor(0x2C, 0x3E, 0x50))

    recs = []
    crm_data = stats.get("crm", {})
    if isinstance(crm_data, dict) and crm_data.get("leads", 0) > 10:
        recs.append(f"CRM has {crm_data['leads']} open leads — schedule follow-up campaigns")
    if isinstance(crm_data, dict) and crm_data.get("opportunities", 0) > 0:
        recs.append(f"{crm_data['opportunities']} opportunities pending — prioritize high-value deals")

    inv_data = stats.get("inventory", {})
    if isinstance(inv_data, dict):
        quants = inv_data.get("stock_quants", 0)
        if quants == 0:
            recs.append("No stock quants found — verify inventory configuration")
        if server_info.get("out_of_stock_count", 0) > 0:
            recs.append(f"{server_info.get('out_of_stock_count')} products out of stock — reorder immediately")

    purch_data = stats.get("purchase", {})
    if isinstance(purch_data, dict) and purch_data.get("draft_rfq", 0) > 0:
        recs.append(f"Purchase: {purch_data['draft_rfq']} draft RFQs — review and send to suppliers")

    acct_data = stats.get("accounting", {})
    if isinstance(acct_data, dict) and acct_data.get("draft", 0) > 0:
        recs.append(f"Accounting: {acct_data['draft']} draft invoices/bills — validate and post")

    exp_data = stats.get("expenses", {})
    if isinstance(exp_data, dict) and exp_data.get("draft", 0) > 0:
        recs.append(f"Expenses: {exp_data['draft']} unreported expenses — remind employees to submit")

    if not recs:
        recs.append("All modules appear healthy — no critical actions needed")

    for i, rec in enumerate(recs):
        add_textbox(slide, 0.8, 1.2 + i * 0.7, 11, 0.6, f"  {i + 1}.  {rec}", size=14, color=RGBColor(0x33, 0x33, 0x33))

    prs.save(str(output_path))
    return str(output_path)


def send_via_gmail(to_email, subject, body, attachments):
    if not GAPI_SCRIPT.exists():
        return {"ok": False, "error": f"Gmail API script not found at {GAPI_SCRIPT}"}

    cmd = [sys.executable, str(GAPI_SCRIPT), "gmail", "send", "--to", to_email, "--subject", subject, "--body", body]

    for att in attachments:
        cmd += ["--attach", att]

    result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
    if result.returncode != 0:
        return {"ok": False, "error": result.stderr.strip()}
    try:
        return json.loads(result.stdout)
    except json.JSONDecodeError:
        return {"ok": True, "result": result.stdout.strip()}


def build_report(args):
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    output_dir = Path(args.output_dir) / f"erp-health-report-{timestamp}"
    output_dir.mkdir(parents=True, exist_ok=True)
    charts_dir = output_dir / "charts"
    charts_dir.mkdir(parents=True, exist_ok=True)

    if args.json:
        print(json.dumps({"status": "started", "output_dir": str(output_dir)}))
        sys.stdout.flush()

    stats = query_module_stats()
    server_info = query_additional_data()

    data_path = output_dir / "data.json"
    data_path.write_text(json.dumps({"stats": stats, "server": server_info}, indent=2, default=str))

    chart_files = generate_charts(stats, output_dir)

    infographic_prompt = generate_infographic_prompt(stats, server_info)
    prompt_path = output_dir / "infographic-prompt.md"
    prompt_path.write_text(infographic_prompt)

    pptx_path = output_dir / "ERP_Health_Report.pptx"
    generate_pptx(stats, server_info, chart_files, pptx_path)

    summary = {
        "status": "ok",
        "timestamp": timestamp,
        "output_dir": str(output_dir),
        "data_file": str(data_path),
        "charts": chart_files,
        "infographic_prompt": str(prompt_path),
        "pptx_report": str(pptx_path),
        "module_stats": {mod: data for mod, data in stats.items() if isinstance(data, dict) and "error" not in data},
        "server_info": server_info.get("server", {}),
    }

    summary_path = output_dir / "summary.json"
    summary_path.write_text(json.dumps(summary, indent=2, default=str))

    return summary


def main():
    parser = argparse.ArgumentParser(description="ERP Health Report — Odoo audit + charts + PPTX + Gmail")
    parser.add_argument("--email", help="Send report via Gmail to this address")
    parser.add_argument("--output-dir", default="./report", help="Output directory for report files")
    parser.add_argument("--json", action="store_true", help="Silent mode, JSON output only")
    parser.add_argument("--skip-email", action="store_true", help="Skip email even if --email is set")
    args = parser.parse_args()

    if not ODOO_SCRIPT.exists():
        msg = json.dumps({"ok": False, "error": f"Odoo script not found at {ODOO_SCRIPT}"})
        print(msg)
        sys.exit(1)

    summary = build_report(args)

    if args.json:
        print(json.dumps(summary, default=str))
    else:
        eprint(f"Report generated: {summary['pptx_report']}")
        eprint(f"Charts: {len(summary['charts'])} files")
        eprint(f"Infographic prompt: {summary['infographic_prompt']}")

    if args.email and not args.skip_email:
        if not GAPI_SCRIPT.exists():
            eprint(f"Warning: Gmail script not found at {GAPI_SCRIPT}. Skipping email.")
            if args.json:
                print(json.dumps({"ok": False, "error": "Gmail not configured"}))
            sys.exit(0)

        attachments = [summary["pptx_report"]]
        combined = summary.get("charts", {}).get("combined")
        if combined and Path(combined).exists():
            attachments.append(combined)

        mod_summary = "; ".join(
            f"{mod}: {sum(v for v in data.values() if isinstance(v, (int, float)))}"
            for mod, data in summary.get("module_stats", {}).items()
        )

        subject = f"ERP Health Report — {datetime.now().strftime('%Y-%m-%d')}"
        body = (
            f"ERP Health Report\n"
            f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}\n\n"
            f"Summary:\n{mod_summary}\n\n"
            f"Attached: PPTX report + infographic overview\n"
            f"---\nAutomated by Hermes ERP Health Report skill"
        )

        result = send_via_gmail(args.email, subject, body, attachments)

        if args.json:
            print(json.dumps({**summary, "email": {"to": args.email, "result": result.get("result", result.get("error", "unknown"))}}, default=str))
        else:
            if result.get("ok"):
                eprint(f"Email sent to {args.email}")
            else:
                eprint(f"Email failed: {result.get('error', 'unknown')}")
    elif args.email and args.skip_email:
        if args.json:
            print(json.dumps({**summary, "email": "skipped"}, default=str))


if __name__ == "__main__":
    main()
